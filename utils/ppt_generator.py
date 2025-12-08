"""
PowerPoint generation utilities
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
import plotly.graph_objects as go
from typing import Dict, List
from config.colours import PIE_COLORS

from utils.chart_generator import generate_channel_airtime_pie  

def hex_to_rgb(hex_color: str) -> tuple:
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def create_single_channel_ppt(config: Dict, 
                              fig_time, 
                              country_chart_img: BytesIO,  
                              stats_text: str,
                              slug_data: Dict) -> BytesIO:
    """
    Generate PowerPoint presentation for single channel with space for images
    
    Args:
        config: Configuration dictionary with channel name, date range, etc.
        fig_time: Time series chart figure
        country_chart_img: Country pie chart image (BytesIO)
        stats_text: Formatted stats text
        slug_data: Slug visualization data
    
    Returns:
        BytesIO object containing PowerPoint file
    """
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Common layout constants (grid)
    LEFT_MARGIN   = Inches(0.5)
    TOP_ROW_Y     = Inches(1.0)
    ROW_GAP       = Inches(0.5)
    BOX_WIDTH     = Inches(3.6)   
    BOX_HEIGHT    = Inches(2.8)   
    COL_GAP       = Inches(0.3)
    
    COL1_X        = LEFT_MARGIN
    COL2_X        = COL1_X + BOX_WIDTH + COL_GAP
    BOTTOM_ROW_Y  = TOP_ROW_Y + BOX_HEIGHT + ROW_GAP

    # Right column placeholder for manual images 
    RIGHT_PLACEHOLDER_WIDTH = Inches(1.6)
    RIGHT_MARGIN            = Inches(0.3)
    RIGHT_PLACEHOLDER_X     = prs.slide_width - RIGHT_PLACEHOLDER_WIDTH - RIGHT_MARGIN

    # Add blank slide
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(7.5), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = f"{config['channel_name']} - Teletrax Analysis ({config.get('date_range', '')})"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    title_para.font.name = 'Arial'
    
    # === Top Left - Time Series Chart (Plotly) ===
    img_bytes = fig_time.to_image(format="png", width=800, height=600)
    img_stream = BytesIO(img_bytes)
    pic = slide.shapes.add_picture(
        img_stream,
        COL1_X, TOP_ROW_Y,
        width=BOX_WIDTH,
        height=BOX_HEIGHT
    )
    pic.line.color.rgb = RGBColor(200, 200, 200)  # Light gray border  
    pic.line.width = Pt(1)

    # === Top Right - 3D Pie Chart (matplotlib BytesIO) ===
    country_chart_img.seek(0)  
    pic2 = slide.shapes.add_picture(
        country_chart_img,
        COL2_X, TOP_ROW_Y,
        width=BOX_WIDTH,
        height=BOX_HEIGHT
    )
    pic2.line.color.rgb = RGBColor(200, 200, 200)  
    pic2.line.width = Pt(1)  
      
    # === Right Column - Placeholder text box for images ===
    placeholder_box = slide.shapes.add_textbox(
        RIGHT_PLACEHOLDER_X, TOP_ROW_Y,
        RIGHT_PLACEHOLDER_WIDTH, Inches(6.0)
    )
    placeholder_frame = placeholder_box.text_frame
    placeholder_frame.text = "[ Add images here ]"
    placeholder_para = placeholder_frame.paragraphs[0]
    placeholder_para.font.size = Pt(10)
    placeholder_para.font.italic = True
    placeholder_para.font.color.rgb = RGBColor(150, 150, 150)
    placeholder_para.alignment = PP_ALIGN.CENTER
    
    # === Bottom Left - Stats Text Box ===
    stats_box = slide.shapes.add_textbox(
        COL1_X, BOTTOM_ROW_Y,
        BOX_WIDTH, BOX_HEIGHT
    )
    # Add border  
    stats_box.line.color.rgb = RGBColor(200, 200, 200)  
    stats_box.line.width = Pt(1)  
    # Add subtle background - can remove if not required
    stats_box.fill.solid()  
    stats_box.fill.fore_color.rgb = RGBColor(250, 250, 250)  # Very light gray

    stats_frame = stats_box.text_frame
    stats_frame.text = stats_text
    stats_frame.word_wrap = True

    for paragraph in stats_frame.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.name = 'Arial'
        paragraph.font.bold = True
        paragraph.space_after = Pt(6)

    
    # === Bottom Right - Top Slug Display ===
    if slug_data:  
        slug_box = slide.shapes.add_textbox(
            COL2_X, BOTTOM_ROW_Y,
            BOX_WIDTH, BOX_HEIGHT
        )  
        # Add border  
        slug_box.line.color.rgb = RGBColor(200, 200, 200)  
        slug_box.line.width = Pt(1)  
        # Add subtle background (optional)
        slug_box.fill.solid()  
        slug_box.fill.fore_color.rgb = RGBColor(250, 250, 250)
        
        slug_frame = slug_box.text_frame  
        slug_frame.word_wrap = True  
        slug_frame.vertical_anchor = 1  # Center vertically
        
        # Large percentage  
        p_pct = slug_frame.paragraphs[0]  
        p_pct.text = f"{slug_data['percentage']:.0f}%"  
        p_pct.font.size = Pt(72)  
        p_pct.font.bold = True  
        p_pct.font.name = 'Arial'  
        p_pct.font.color.rgb = RGBColor(217, 120, 71)  # Orange color  
        p_pct.alignment = PP_ALIGN.CENTER  
        p_pct.space_after = Pt(8)
        
        # Text: "of detections since [date] about [slug]"  
        p_text = slug_frame.add_paragraph()  
        earliest_date = slug_data.get('earliest_date', '2024')  
        p_text.text = f"of detections since {earliest_date}\nabout {slug_data['clean_slug']}"  
        p_text.font.size = Pt(13)  
        p_text.font.name = 'Arial'  
        p_text.font.bold = False  
        p_text.font.color.rgb = RGBColor(100, 100, 100)  
        p_text.alignment = PP_ALIGN.CENTER   

    # Save to BytesIO
    ppt_stream = BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    
    return ppt_stream


def create_multi_channel_ppt(config: Dict,
                             fig_time,
                             country_chart_img: BytesIO,
                             stats_text: str,
                             slug_data: Dict,
                             channel_charts: Dict[str, BytesIO],
                             comprehensive_data=None) -> BytesIO:
    """
    Generate PowerPoint presentation for multiple channels with space for images
    
    Args:
        config: Configuration dictionary
        fig_time: Time series Plotly chart figure
        country_chart_img: Main country pie chart (BytesIO matplotlib image)
        stats_text: Formatted stats text
        slug_data: Slug visualization data
        channel_charts: Dictionary of channel-specific charts (BytesIO images)
    
    Returns:
        BytesIO object containing PowerPoint file
    """
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Shared layout constants (reuse same grid as single-channel)
    LEFT_MARGIN   = Inches(0.5)
    TOP_ROW_Y     = Inches(1.0)
    ROW_GAP       = Inches(0.5)
    BOX_WIDTH     = Inches(3.6)
    BOX_HEIGHT    = Inches(2.8)
    COL_GAP       = Inches(0.3)

    COL1_X        = LEFT_MARGIN
    COL2_X        = COL1_X + BOX_WIDTH + COL_GAP
    BOTTOM_ROW_Y  = TOP_ROW_Y + BOX_HEIGHT + ROW_GAP

    RIGHT_PLACEHOLDER_WIDTH = Inches(1.6)
    RIGHT_MARGIN            = Inches(0.3)
    RIGHT_PLACEHOLDER_X     = prs.slide_width - RIGHT_PLACEHOLDER_WIDTH - RIGHT_MARGIN
    
    # === SLIDE 1: Main Overview ===
    blank_slide_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_slide_layout)
    
    # Title
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(7.5), Inches(0.5))
    title_frame = title_box.text_frame
    channel_names = config.get('channel_names', [config['channel_name']])
    title_frame.text = f"{', '.join(channel_names)} - Teletrax Analysis ({config.get('date_range', '')})"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(22)
    title_para.font.bold = True
    title_para.font.name = 'Arial'
    
    # Top Left - Time Series (Plotly)
    img_bytes = fig_time.to_image(format="png", width=800, height=600)
    img_stream = BytesIO(img_bytes)
    slide1.shapes.add_picture(
        img_stream,
        COL1_X, TOP_ROW_Y,
        width=BOX_WIDTH,
        height=BOX_HEIGHT
    )
    
    # Top Right - Channel Airtime Distribution (for multi-channel)  
    if comprehensive_data is not None:  
        try:  
            # Generate channel airtime chart  
            channel_airtime_img = generate_channel_airtime_pie(  
                comprehensive_data,  
                title="Time on air by channel",  
                subtitle="Material under 30 days"  
            )
            
            # Check if chart was generated successfully  
            if channel_airtime_img.getbuffer().nbytes > 0:  
                channel_airtime_img.seek(0)  
                pic2 = slide1.shapes.add_picture(  
                    channel_airtime_img,  
                    COL2_X, TOP_ROW_Y,  
                    width=BOX_WIDTH,  
                    height=BOX_HEIGHT  
                )  
                pic2.line.color.rgb = RGBColor(200, 200, 200)  
                pic2.line.width = Pt(1)  
            else:  
                # Fallback to country chart  
                country_chart_img.seek(0)  
                pic2 = slide1.shapes.add_picture(  
                    country_chart_img,  
                    COL2_X, TOP_ROW_Y,  
                    width=BOX_WIDTH,  
                    height=BOX_HEIGHT  
                )  
                pic2.line.color.rgb = RGBColor(200, 200, 200)  
                pic2.line.width = Pt(1)  
        except Exception as e:  
            # If any error, fallback to country chart  
            country_chart_img.seek(0)  
            pic2 = slide1.shapes.add_picture(  
                country_chart_img,  
                COL2_X, TOP_ROW_Y,  
                width=BOX_WIDTH,  
                height=BOX_HEIGHT  
            )  
            pic2.line.color.rgb = RGBColor(200, 200, 200)  
            pic2.line.width = Pt(1)  
    else:  
        # No comprehensive data provided - use country chart  
        country_chart_img.seek(0)  
        pic2 = slide1.shapes.add_picture(  
            country_chart_img,  
            COL2_X, TOP_ROW_Y,  
            width=BOX_WIDTH,  
            height=BOX_HEIGHT  
        )  
        pic2.line.color.rgb = RGBColor(200, 200, 200)  
        pic2.line.width = Pt(1)
    
    # Placeholder for images on the right   
    placeholder_box = slide1.shapes.add_textbox(  
        RIGHT_PLACEHOLDER_X, TOP_ROW_Y,  
        RIGHT_PLACEHOLDER_WIDTH, Inches(6.0)  
    )  
    placeholder_frame = placeholder_box.text_frame  
    placeholder_frame.text = "[ Add images here ]"  
    placeholder_para = placeholder_frame.paragraphs[0]  
    placeholder_para.font.size = Pt(10)  
    placeholder_para.font.italic = True  
    placeholder_para.font.color.rgb = RGBColor(150, 150, 150)  
    placeholder_para.alignment = PP_ALIGN.CENTER  
    
    # Bottom Left - Stats
    stats_box = slide1.shapes.add_textbox(
        COL1_X, BOTTOM_ROW_Y,
        BOX_WIDTH, BOX_HEIGHT
    )
    stats_frame = stats_box.text_frame
    stats_frame.text = stats_text
    stats_frame.word_wrap = True
    
    for paragraph in stats_frame.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.name = 'Arial'
        paragraph.font.bold = True
        paragraph.space_after = Pt(6)
    
    # Bottom Right - Top Slug
    if slug_data:  
        slug_box = slide1.shapes.add_textbox(
            COL2_X, BOTTOM_ROW_Y,
            BOX_WIDTH, BOX_HEIGHT
        )  
        # Add border  
        slug_box.line.color.rgb = RGBColor(200, 200, 200)  
        slug_box.line.width = Pt(1)  
        # Optional: add subtle background  
        slug_box.fill.solid()  
        slug_box.fill.fore_color.rgb = RGBColor(250, 250, 250)
        
        slug_frame = slug_box.text_frame  
        slug_frame.word_wrap = True  
        slug_frame.vertical_anchor = 1  # Center vertically
        
        # Large percentage in ORANGE  
        p_pct = slug_frame.paragraphs[0]  
        p_pct.text = f"{slug_data['percentage']:.0f}%"  
        p_pct.font.size = Pt(72)  
        p_pct.font.bold = True  
        p_pct.font.name = 'Arial'  
        p_pct.font.color.rgb = RGBColor(217, 120, 71)  # Orange color  
        p_pct.alignment = PP_ALIGN.CENTER  
        p_pct.space_after = Pt(8)
        
        # Text: "of detections since [date] about [slug]"  
        p_text = slug_frame.add_paragraph()  
        earliest_date = slug_data.get('earliest_date', '2024')  
        p_text.text = f"of detections since {earliest_date}\nabout {slug_data['clean_slug']}"  
        p_text.font.size = Pt(13)  
        p_text.font.name = 'Arial'  
        p_text.font.bold = False  
        p_text.font.color.rgb = RGBColor(100, 100, 100)  
        p_text.alignment = PP_ALIGN.CENTER   

    # === SLIDE 2: Channel Breakdown ===
    if channel_charts:
        slide2 = prs.slides.add_slide(blank_slide_layout)
        
        # Title
        title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(7.5), Inches(0.5))
        title_frame2 = title_box2.text_frame
        title_frame2.text = "Usage by Channel - Asset Age <30 Days"
        title_para2 = title_frame2.paragraphs[0]
        title_para2.font.size = Pt(22)
        title_para2.font.bold = True
        title_para2.font.name = 'Arial'
        
        # Add channel charts in grid layout 
        col_width = 3.8
        col_spacing = 0.4
        row_height = 3.0
        row_spacing = 0.3
        start_x = 0.5
        start_y = 1.0
        
        for idx, (channel_name, img_buffer) in enumerate(channel_charts.items()):
            col = idx % 2
            row = idx // 2
            
            x_pos = start_x + (col * (col_width + col_spacing))
            y_pos = start_y + (row * (row_height + row_spacing))
            
            # Add matplotlib image from BytesIO
            img_buffer.seek(0)
            slide2.shapes.add_picture(img_buffer, Inches(x_pos), Inches(y_pos), width=Inches(col_width))
        
        # Placeholder on slide 2 as well
        placeholder_box2 = slide2.shapes.add_textbox(
            RIGHT_PLACEHOLDER_X, TOP_ROW_Y,
            RIGHT_PLACEHOLDER_WIDTH, Inches(6.0)
        )
        placeholder_frame2 = placeholder_box2.text_frame
        placeholder_frame2.text = "[ Add images here ]"
        placeholder_para2 = placeholder_frame2.paragraphs[0]
        placeholder_para2.font.size = Pt(10)
        placeholder_para2.font.italic = True
        placeholder_para2.font.color.rgb = RGBColor(150, 150, 150)
        placeholder_para2.alignment = PP_ALIGN.CENTER
    
    # Save to BytesIO
    ppt_stream = BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    
    return ppt_stream